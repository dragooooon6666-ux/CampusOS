"""
文件内容提取模块

支持：docx / pdf / xlsx / pptx / txt / md / 纯文本
"""

import re
from pathlib import Path

# ── 提取结果 ──────────────────────────────────────────────

MAX_PREVIEW_CHARS = 800
TITLE_LINES = 6
SIGNATURE_LINES = 4


class ContentPreview:
    """提取后的内容预览，包含标题区和落款区"""

    def __init__(self, title_section: str = "", signature_section: str = ""):
        self.title_section = title_section
        self.signature_section = signature_section
        self.full_text = ""

    @property
    def ai_context(self) -> str:
        parts = []
        if self.title_section:
            parts.append(f"【文件开头内容】\n{self.title_section}")
        if self.signature_section:
            parts.append(f"【文件末尾内容（落款/署名）】\n{self.signature_section}")
        return "\n\n".join(parts) if parts else ""

    def __bool__(self):
        return bool(self.ai_context)


# ── 各类型提取器 ───────────────────────────────────────────

def _extract_docx(file_path: Path) -> ContentPreview:
    import docx
    doc = docx.Document(str(file_path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return ContentPreview()
    return _build_preview(paragraphs)


def _extract_pdf(file_path: Path) -> ContentPreview:
    from PyPDF2 import PdfReader
    reader = PdfReader(str(file_path))
    all_lines: list[str] = []
    for page in reader.pages[:3]:
        text = page.extract_text()
        if text:
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            all_lines.extend(lines)
        if sum(len(l) for l in all_lines) > MAX_PREVIEW_CHARS * 2:
            break
    if not all_lines:
        return ContentPreview()
    return _build_preview(all_lines)


def _extract_xlsx(file_path: Path) -> ContentPreview:
    from openpyxl import load_workbook
    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet_name in wb.sheetnames[:2]:
        ws = wb[sheet_name]
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
            if row_idx > 20:
                break
            row_text = " ".join(str(c) for c in row if c is not None).strip()
            if row_text:
                lines.append(row_text)
    wb.close()
    return _build_preview(lines)


def _extract_pptx(file_path: Path) -> ContentPreview:
    """从 PowerPoint 提取幻灯片文本"""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    all_lines: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        all_lines.append(text)
        if sum(len(l) for l in all_lines) > MAX_PREVIEW_CHARS * 2:
            break

    if not all_lines:
        return ContentPreview()
    return _build_preview(all_lines)


def _extract_plain(file_path: Path) -> ContentPreview:
    try:
        text = file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            text = file_path.read_text(encoding="gbk")
        except UnicodeDecodeError:
            return ContentPreview()
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return _build_preview(lines)


def _extract_nothing(file_path: Path) -> ContentPreview:
    return ContentPreview()


# ── 通用构建 ───────────────────────────────────────────────

def _build_preview(lines: list[str]) -> ContentPreview:
    if not lines:
        return ContentPreview()

    preview = ContentPreview()

    title_chars = 0
    title_lines = []
    for line in lines[:TITLE_LINES]:
        title_lines.append(line)
        title_chars += len(line)
        if title_chars >= MAX_PREVIEW_CHARS // 2:
            break
    preview.title_section = "\n".join(title_lines)[:MAX_PREVIEW_CHARS]

    if len(lines) > TITLE_LINES + 2:
        sig = lines[-SIGNATURE_LINES:]
        preview.signature_section = "\n".join(sig)[:MAX_PREVIEW_CHARS // 2]

    preview.full_text = "\n".join(lines)
    return preview


# ── 文件名关键词检测 ────────────────────────────────────────

# 文件名中包含这些关键词时，AI 分类需优先考虑
FILENAME_TYPE_HINTS = {
    # 文稿类
    "策划": "策划案", "策划书": "策划案", "策划案": "策划案",
    "新闻稿": "新闻稿", "新闻": "新闻稿", "通讯": "新闻稿",
    "总结": "活动总结", "活动总结": "活动总结",
    "纪要": "会议纪要", "会议纪要": "会议纪要", "会议记录": "会议纪要",
    "通知": "通知", "公告": "通知",
    "请示": "请示", "申请": "申请书", "申请书": "申请书", "申报表": "申请书", "申报": "申请书",
    "发言稿": "发言稿", "讲话": "发言稿", "致辞": "发言稿", "主持稿": "发言稿",
    "汇报": "报告", "报告": "报告",
    "述职": "述职报告", "述职报告": "述职报告",
    "简介": "简历", "简历": "简历", "个人简历": "简历",
    "证明": "证明", "复印件": "证明", "证书": "证明",
    "方案": "方案", "指南": "方案", "提纲": "方案",
    "议程": "议程", "班会": "议程",
    "作业": "作业",
    "论文": "论文", "毕业论文": "论文",
    "笔记": "笔记", "记录": "笔记",
    # 表格类
    "预算": "预算表", "预算表": "预算表", "经费": "预算表",
    "签到": "签到表", "签到表": "签到表",
    "物资": "物资清单", "清单": "物资清单",
    "统计": "统计表", "统计表": "统计表", "汇总": "统计表",
    "通讯录": "通讯录", "联系方式": "通讯录",
    "排班": "排班表", "值班": "排班表",
    # 收集表类
    "收集表": "收集表",
    "报名表": "收集表",
    "异动表": "收集表",
    "反馈表": "收集表",
    "登记表": "收集表",
    "申报表": "收集表",
    "审批表": "收集表",
    "信息表": "收集表",
    "调查表": "收集表",
    # 统计表
    "名单": "统计表",
    "一览表": "统计表",
}


def detect_type_from_filename(filename: str) -> str | None:
    """从文件名检测文档类型，返回类型或 None。长关键词优先匹配。"""
    stem = Path(filename).stem
    # 按关键词长度降序排列，避免短词误匹配（如"通讯"不应匹配"通讯录"）
    sorted_hints = sorted(FILENAME_TYPE_HINTS.items(), key=lambda x: len(x[0]), reverse=True)
    for keyword, doc_type in sorted_hints:
        if keyword in stem:
            return doc_type
    return None


# ── 类型 → 提取器映射 ──────────────────────────────────────

EXTRACTOR_MAP = {
    ".docx": _extract_docx,
    ".doc": _extract_docx,
    ".pdf": _extract_pdf,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xlsx,
    ".csv": _extract_plain,
    ".pptx": _extract_pptx,
    ".ppt": _extract_pptx,
    ".txt": _extract_plain,
    ".md": _extract_plain,
    ".markdown": _extract_plain,
    ".log": _extract_plain,
}

FALLBACK_EXTRACTOR = _extract_nothing


# ── 公开 API ───────────────────────────────────────────────

def extract_content_preview(file_path: Path) -> ContentPreview:
    """提取文件内容预览，包含标题区和落款区。"""
    extension = file_path.suffix.lower()
    extractor = EXTRACTOR_MAP.get(extension, FALLBACK_EXTRACTOR)

    try:
        preview = extractor(file_path)
    except Exception:
        preview = ContentPreview()

    # 文件名作为强信号 fallback
    if not preview.ai_context:
        stem = file_path.stem.strip()
        if stem and len(stem) >= 2:
            preview.title_section = f"（文件名为：{stem}）"

    return preview
